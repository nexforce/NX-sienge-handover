"""
Nexforce brand helpers for python-pptx.

Usage:
    from pptx import Presentation
    from pptx.util import Cm
    from python_pptx_helpers import (
        apply_nexforce_theme,
        cover_slide,
        content_slide,
        section_divider_slide,
        closing_slide,
        add_brand_table,
        add_kpi_block,
        NEXFORCE_COLORS,
        logo_path,
    )

    prs = Presentation()
    apply_nexforce_theme(prs)

    cover_slide(prs, title="Q1 Results", subtitle="Jan-Mar 2026",
                meta="Confidencial - Uso interno")

    s = content_slide(prs, title="Highlights")
    add_kpi_block(s, Cm(2), Cm(5),
                   value="18%", label="Growth YoY",
                   delta="+3pp", delta_kind="positive")

    closing_slide(prs, message="Obrigado.", contact="team@nexforce.ai")
    prs.save("deck.pptx")

All functions enforce Nexforce brand rules: Lato font, palette-only colors,
semantic colors never on slide or cell backgrounds except the defined black
header/total rows, and logos resolved from the skill's local assets folder.
"""

from __future__ import annotations

import glob
import struct
from functools import lru_cache
from pathlib import Path
from typing import Optional, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as PresType
from pptx.slide import Slide
from pptx.util import Cm, Emu, Pt


# === Palette ==============================================================

NEXFORCE_COLORS: dict[str, RGBColor] = {
    "primary_black": RGBColor(0x0C, 0x0E, 0x0E),
    "white":         RGBColor(0xFF, 0xFF, 0xFF),
    "dark_gray":     RGBColor(0x51, 0x51, 0x51),
    "mid_gray":      RGBColor(0x77, 0x77, 0x77),
    "light_gray":    RGBColor(0x9C, 0x9B, 0x9B),
    "near_white":    RGBColor(0xF5, 0xF5, 0xF5),
    "navy":          RGBColor(0x30, 0x3F, 0x63),
    "blue":          RGBColor(0x21, 0x5A, 0x9F),
    "green":         RGBColor(0x2D, 0x6E, 0x44),
    "yellow":        RGBColor(0xD8, 0xB5, 0x23),
    "red":           RGBColor(0xBA, 0x19, 0x25),
    "orange":        RGBColor(0xD5, 0x63, 0x16),
    "purple":        RGBColor(0x66, 0x2A, 0x7F),
    "pink":          RGBColor(0xB1, 0x33, 0x8A),
}


# === Asset paths ==========================================================

def nexforce_assets() -> Path:
    """Locate the nexforce-brand/assets folder across environments."""
    candidates = [
        # 1) claude.ai computer-use (bash_tool sandbox)
        Path("/mnt/skills/user/nexforce-brand/assets"),
        # 2) Cowork Linux sandbox
        *[Path(p) for p in glob.glob("/sessions/*/mnt/Work/Skills/nexforce-brand/assets")],
        *[Path(p) for p in glob.glob("/sessions/*/mnt/.claude/skills/nexforce-brand/assets")],
        # 3) Windows local
        Path.home() / "Desktop" / "Work" / "Skills" / "nexforce-brand" / "assets",
    ]
    for p in candidates:
        if p.exists():
            return p
    raise FileNotFoundError(
        "nexforce-brand/assets folder not found. Check skill installation."
    )


def logo_path(variant: str) -> str:
    """Absolute path to logo-black.png or logo-white.png."""
    if variant not in ("black", "white"):
        raise ValueError("variant must be 'black' or 'white'")
    return str(nexforce_assets() / f"logo-{variant}.png")


@lru_cache(maxsize=4)
def logo_aspect(variant: str) -> float:
    """Width/height aspect ratio of the logo PNG.

    Reads the PNG IHDR chunk directly, no Pillow dependency.
    The Nexforce wordmark is ~5.54:1, so a 2cm tall logo is ~11cm wide --
    using a square assumption silently breaks layouts.
    """
    path = logo_path(variant)
    with open(path, "rb") as f:
        sig = f.read(8)
        if sig[:8] != b"\x89PNG\r\n\x1a\n":
            raise ValueError(f"{path} is not a PNG file")
        f.read(8)  # length + "IHDR"
        width, height = struct.unpack(">II", f.read(8))
    return width / height


def logo_width_for_height(variant: str, height_emu: Emu) -> int:
    """Compute the rendered width (EMU) for a target logo height."""
    return int(round(height_emu * logo_aspect(variant)))


# === Theme ================================================================

def apply_nexforce_theme(prs: PresType) -> None:
    """Configure 16:9 widescreen canvas."""
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)


# === Slide primitives =====================================================

def _blank_slide(prs: PresType) -> Slide:
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _fill_slide_background(slide: Slide, color_key: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NEXFORCE_COLORS[color_key]


def _add_logo(slide: Slide, variant: str, left: Emu, top: Emu,
               height: Emu) -> None:
    slide.shapes.add_picture(logo_path(variant), left, top, height=height)


def _add_text(slide: Slide, left: Emu, top: Emu, width: Emu, height: Emu, *,
               text: str, size: int, color_key: str,
               bold: bool = False, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Lato"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = NEXFORCE_COLORS[color_key]


# === Slide builders =======================================================

def cover_slide(prs: PresType, title: str,
                 subtitle: Optional[str] = None,
                 meta: Optional[str] = None) -> Slide:
    """Dark cover slide with white logo top-left and large title."""
    slide = _blank_slide(prs)
    _fill_slide_background(slide, "primary_black")
    # Wordmark ~5.54:1 -- height 1.6cm renders ~8.9cm wide.
    _add_logo(slide, "white", Cm(2), Cm(1.5), Cm(1.6))

    _add_text(slide, Cm(2), Cm(8), Cm(28), Cm(4),
              text=title, size=44, color_key="white", bold=True)

    if subtitle:
        _add_text(slide, Cm(2), Cm(11.5), Cm(28), Cm(2),
                  text=subtitle, size=20, color_key="light_gray")

    if meta:
        _add_text(slide, Cm(2), Cm(17), Cm(28), Cm(1),
                  text=meta, size=10, color_key="light_gray")

    return slide


def content_slide(prs: PresType, title: str) -> Slide:
    """White content slide with black logo top-left and navy title."""
    slide = _blank_slide(prs)
    _fill_slide_background(slide, "white")
    # Subtle header logo -- 0.8cm tall renders ~4.4cm wide.
    _add_logo(slide, "black", Cm(1.5), Cm(0.8), Cm(0.8))

    _add_text(slide, Cm(1.5), Cm(2.8), Cm(30), Cm(2),
              text=title, size=28, color_key="navy", bold=True)

    return slide


def section_divider_slide(prs: PresType, title: str) -> Slide:
    """Full dark slide with large white title -- use between sections."""
    slide = _blank_slide(prs)
    _fill_slide_background(slide, "primary_black")

    _add_text(slide, Cm(2), Cm(8), Cm(30), Cm(4),
              text=title, size=54, color_key="white", bold=True)

    return slide


def closing_slide(prs: PresType, message: str = "Obrigado.",
                   contact: Optional[str] = None) -> Slide:
    """Dark closing slide with centered logo and optional contact line.

    Logo width is computed from the PNG aspect ratio (~5.54:1 wordmark),
    then centered against the actual slide width.
    """
    slide = _blank_slide(prs)
    _fill_slide_background(slide, "primary_black")

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    logo_h = Cm(2.2)
    logo_w = logo_width_for_height("white", logo_h)
    logo_left = (slide_w - logo_w) // 2
    logo_top = int(slide_h * 0.32)
    _add_logo(slide, "white", logo_left, logo_top, logo_h)

    msg_top = logo_top + logo_h + Cm(1.5)
    _add_text(slide, Cm(2), msg_top, slide_w - Cm(4), Cm(2),
              text=message, size=32, color_key="white",
              align=PP_ALIGN.CENTER)

    if contact:
        _add_text(slide, Cm(2), slide_h - Cm(2), slide_w - Cm(4), Cm(1),
                  text=contact, size=10, color_key="light_gray",
                  align=PP_ALIGN.CENTER)

    return slide


# === Table builder ========================================================

def add_brand_table(slide: Slide, left: Emu, top: Emu,
                     width: Emu, height: Emu,
                     headers: List[str],
                     rows: List[List[str]]) -> None:
    """Branded table: black header row, alternating #FFFFFF / #F5F5F5 body."""
    n_cols = len(headers)
    n_rows = 1 + len(rows)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table

    for i, text in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NEXFORCE_COLORS["primary_black"]
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = str(text)
        r.font.name = "Lato"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = NEXFORCE_COLORS["white"]

    for r_idx, row_values in enumerate(rows):
        bg = (NEXFORCE_COLORS["white"] if r_idx % 2 == 0
              else NEXFORCE_COLORS["near_white"])
        for c_idx, value in enumerate(row_values):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            rr = p.add_run()
            rr.text = str(value)
            rr.font.name = "Lato"
            rr.font.size = Pt(10)
            rr.font.color.rgb = NEXFORCE_COLORS["dark_gray"]


# === KPI block ============================================================

_DELTA_KINDS = {
    "positive":  "green",
    "attention": "yellow",
    "alert":     "red",
}


def add_kpi_block(slide: Slide, left: Emu, top: Emu, *,
                   value: str, label: str,
                   delta: Optional[str] = None,
                   delta_kind: Optional[str] = None) -> None:
    """Large KPI number, optional inline colored delta, caption below."""
    tb = slide.shapes.add_textbox(left, top, Cm(10), Cm(4))
    tf = tb.text_frame

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.name = "Lato"
    r.font.size = Pt(48)
    r.font.bold = True
    r.font.color.rgb = NEXFORCE_COLORS["primary_black"]

    if delta and delta_kind:
        if delta_kind not in _DELTA_KINDS:
            raise ValueError(
                f"Unknown delta_kind {delta_kind!r}; "
                f"allowed: {sorted(_DELTA_KINDS.keys())}"
            )
        dr = p.add_run()
        dr.text = "  " + delta
        dr.font.name = "Lato"
        dr.font.size = Pt(16)
        dr.font.bold = True
        dr.font.color.rgb = NEXFORCE_COLORS[_DELTA_KINDS[delta_kind]]

    lp = tf.add_paragraph()
    lp.space_before = Pt(8)
    lr = lp.add_run()
    lr.text = label
    lr.font.name = "Lato"
    lr.font.size = Pt(12)
    lr.font.color.rgb = NEXFORCE_COLORS["mid_gray"]


# === Footer ===============================================================

def icon_mark_path(on_dark: bool = False) -> str:
    """Absolute path to the icon-mark PNG variant.

    icon-black.png: dark badge with white X -- use on LIGHT backgrounds.
    icon-white.png: light badge with dark X -- use on DARK backgrounds.
    """
    variant = "white" if on_dark else "black"
    return str(nexforce_assets() / f"icon-{variant}.png")


def add_brand_footer(slide: Slide, prs: PresType, *,
                     meta: str = "",
                     on_dark: bool = False,
                     icon_size: Emu = Cm(0.6)) -> None:
    """Standard Nexforce footer: icon-mark bottom-left, meta bottom-right.

    The icon-mark is used alone -- no 'Nexforce' wordmark text.

    Args:
        slide: Slide to decorate.
        prs: Presentation (for slide_width / slide_height).
        meta: Optional right-aligned metadata string (page number, confidentiality,
              date, etc.). Empty string draws only the icon.
        on_dark: Set True for slides with dark background so the light-badge
                 variant is used.
        icon_size: Icon height (square badge). Default 0.6 cm.
    """
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    margin = Cm(1.2)
    icon_left = margin
    icon_top = slide_h - margin - icon_size + Cm(0.2)
    slide.shapes.add_picture(icon_mark_path(on_dark),
                              icon_left, icon_top, height=icon_size)

    if meta:
        meta_color = "light_gray" if on_dark else "mid_gray"
        text_width = Cm(20)
        text_left = slide_w - margin - text_width
        text_top = icon_top + Cm(0.05)
        _add_text(slide, text_left, text_top, text_width, icon_size,
                  text=meta, size=9, color_key=meta_color,
                  align=PP_ALIGN.RIGHT)
